"""Agent loop for the Al Dente Company Brain. v3

Tools:
  call_api              – fully-paginated calls to ALL Al Dente mock API endpoints
  search_kb             – keyword + SKU-index retrieval over backend/data/kb/
  search_transcripts    – batch-scan all call transcripts for a keyword (Python loop)
  crm_channel_breakdown – Python join: opportunities × customers → totals by channel

All arithmetic (sums, counts, group-bys) is computed in Python, never in the LLM.
Pagination: every list endpoint is fully consumed before the result is returned.
"""

from __future__ import annotations

import io
import json
import os
import re
import time
import uuid
from pathlib import Path

import httpx
from openai import OpenAI

# ── Static files directory (binary artifacts) ────────────────────────────────
_FILES_DIR = Path(__file__).parent / "static" / "files"
_FILES_DIR.mkdir(parents=True, exist_ok=True)


# ── Knowledge-base ────────────────────────────────────────────────────────────

_KB_DIR = Path(__file__).parent / "data" / "kb"
_KB: dict[str, str] = {
    f.stem: f.read_text(encoding="utf-8") for f in sorted(_KB_DIR.glob("*.md"))
}

# Reverse index: finished-good/raw-material SKU → first KB doc that mentions it
_SKU_RE = re.compile(r"\b(PAS-[A-Z]{3}-\d{3}|RAW-[A-Z]{3}-\d{3})\b")
_KB_SKU_INDEX: dict[str, list[str]] = {}   # "PAS-SPA-500" → ["DOC-001", ...]
for _doc_id, _content in _KB.items():
    for _sku in _SKU_RE.findall(_content):
        _KB_SKU_INDEX.setdefault(_sku, [])
        if _doc_id not in _KB_SKU_INDEX[_sku]:
            _KB_SKU_INDEX[_sku].append(_doc_id)


def _search_kb_impl(query: str, doc_ids: list[str] | None = None) -> list[dict]:
    """Whole-document retrieval. Prioritises exact SKU matches, then keyword score."""
    if doc_ids:
        return [{"id": d, "content": _KB[d]} for d in doc_ids if d in _KB]

    results: list[dict] = []
    seen: set[str] = set()

    # 1. Exact SKU hit (highest priority)
    skus_in_query = _SKU_RE.findall(query.upper())
    for sku in skus_in_query:
        for did in _KB_SKU_INDEX.get(sku, []):
            if did not in seen:
                results.append({"id": did, "content": _KB[did]})
                seen.add(did)

    # 2. Keyword relevance for remaining slots (up to 5 total)
    q_lower = query.lower()
    q_words = {w for w in q_lower.split() if len(w) > 2}
    scored: list[tuple[int, str]] = []
    for doc_id, content in _KB.items():
        if doc_id in seen:
            continue
        c_lower = content.lower()
        score = sum(1 for w in q_words if w in c_lower)
        if q_lower in c_lower:
            score += 10
        if score > 0:
            scored.append((score, doc_id))
    scored.sort(reverse=True)
    for _, did in scored[: max(0, 5 - len(results))]:
        results.append({"id": did, "content": _KB[did]})

    return results


# ── Binary artifact generators ───────────────────────────────────────────────

def _make_pdf(title: str, sections: list[dict]) -> bytes:
    from fpdf import FPDF  # type: ignore

    # A4 portrait: 210mm wide, 10mm margins each side → 190mm usable
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_margins(10, 10, 10)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    eff_w = pdf.w - pdf.l_margin - pdf.r_margin  # 190.0 mm

    # ── Title ───────────────────────────────────────────────────────
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(eff_w, 10, title[:80], new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    for sec in sections:
        # Section heading
        if heading := sec.get("heading", ""):
            pdf.set_font("Helvetica", "B", 12)
            pdf.cell(eff_w, 8, heading[:80], new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

        # Free text block
        if text := sec.get("text", ""):
            pdf.set_font("Helvetica", "", 10)
            safe = text.replace("\r\n", "\n").replace("\r", "\n")
            pdf.multi_cell(eff_w, 6, safe[:2000])
            pdf.ln(2)

        # Table rows
        rows = sec.get("rows") or []
        if not rows:
            continue

        n_cols = len(rows[0])

        # Column widths that sum exactly to eff_w (190 mm)
        if n_cols == 4:
            col_widths = [38.0, 82.0, 35.0, 35.0]   # SKU | Desc | On-Hand | Min
        elif n_cols == 3:
            col_widths = [40.0, 110.0, 40.0]
        elif n_cols == 2:
            col_widths = [85.0, 105.0]
        elif n_cols == 1:
            col_widths = [eff_w]
        else:
            w = round(eff_w / n_cols, 2)
            col_widths = [w] * n_cols
            col_widths[-1] = round(eff_w - w * (n_cols - 1), 2)  # absorb rounding

        ROW_H = 7.0

        for ri, row in enumerate(rows):
            is_header = (ri == 0)
            if is_header:
                pdf.set_font("Helvetica", "B", 10)
                pdf.set_fill_color(220, 220, 220)
            else:
                pdf.set_font("Helvetica", "", 9)

            pdf.set_x(pdf.l_margin)
            for ci, (cell_val, cw) in enumerate(zip(row, col_widths)):
                raw = str(cell_val)
                # Estimate max chars that fit: Helvetica ~1.85mm/char at 9-10pt
                pt = 10 if is_header else 9
                max_chars = max(4, int(cw / (pt * 0.19)))
                if len(raw) > max_chars:
                    raw = raw[:max_chars - 1] + "…"
                align = "C" if is_header else ("R" if ci >= 2 else "L")
                pdf.cell(cw, ROW_H, raw, border=1, align=align,
                         fill=is_header, new_x="RIGHT", new_y="TOP")
            pdf.ln(ROW_H)

        pdf.ln(3)

    # Strip trailing blank pages: fpdf2 may produce one if auto-page-break
    # fired right at the end. We can't easily detect this without hacks, but
    # by NOT calling add_page() after the last section we avoid it.
    return bytes(pdf.output())


def _make_docx(title: str, sections: list[dict]) -> bytes:
    from docx import Document  # type: ignore
    from docx.shared import Pt  # type: ignore

    doc = Document()
    doc.add_heading(title, level=0)
    for sec in sections:
        doc.add_heading(sec.get("heading", ""), level=1)
        if text := sec.get("text", ""):
            doc.add_paragraph(text)
        if rows := sec.get("rows", []):
            if rows:
                tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
                tbl.style = "Table Grid"
                for ri, row in enumerate(rows):
                    for ci, cell in enumerate(row):
                        tbl.cell(ri, ci).text = str(cell)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(title: str, sections: list[dict]) -> bytes:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    prs = Presentation()
    # Title slide
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = title
    if ts.placeholders[1]:
        ts.placeholders[1].text = "Al Dente S.r.l."
    # Content slides
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sec.get("heading", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        lines: list[str] = []
        if text := sec.get("text", ""):
            lines += [t.strip() for t in text.split("\n") if t.strip()]
        if rows := sec.get("rows", []):
            for row in rows:
                lines.append("  |  ".join(str(c) for c in row))
        for i, line in enumerate(lines[:12]):  # cap at 12 bullets per slide
            if i == 0:
                tf.text = line
            else:
                tf.add_paragraph().text = line
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(title: str, sections: list[dict]) -> bytes:
    from openpyxl import Workbook  # type: ignore
    from openpyxl.styles import Font  # type: ignore

    wb = Workbook()
    default_sheet = wb.active
    for sec in sections:
        name = sec.get("heading", "Sheet")[:31] or "Sheet"
        ws = wb.create_sheet(title=name)
        ws.append([title])
        ws["A1"].font = Font(bold=True, size=13)
        ws.append([])
        if text := sec.get("text", ""):
            for line in text.split("\n"):
                ws.append([line])
            ws.append([])
        if rows := sec.get("rows", []):
            for ri, row in enumerate(rows):
                ws.append(list(row))
                if ri == 0:  # bold header row
                    for cell in ws[ws.max_row]:
                        cell.font = Font(bold=True)
    wb.remove(default_sheet)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


_FORMAT_GENERATORS = {
    "pdf": (_make_pdf, ".pdf"),
    "docx": (_make_docx, ".docx"),
    "pptx": (_make_pptx, ".pptx"),
    "xlsx": (_make_xlsx, ".xlsx"),
}


def _generate_artifact_impl(fmt: str, title: str, sections: list[dict]) -> dict:
    """
    Sections may include:
      - 'rows': explicit table data (first row = headers)
      - 'data_source': API endpoint string; Python fetches + auto-formats it
        optionally combined with 'column_fields' (list of field names to extract)
        and 'column_headers' (display header labels)
    """
    fmt = fmt.lower().strip()
    if fmt not in _FORMAT_GENERATORS:
        return {"error": f"unsupported format: {fmt}. Use pdf, docx, pptx, or xlsx."}
    try:
        # Resolve any data_source entries into rows
        for sec in sections:
            if "data_source" in sec and not sec.get("rows"):
                endpoint = sec.pop("data_source")
                col_fields: list[str] = sec.pop("column_fields", [])
                col_headers: list[str] = sec.pop("column_headers", [])
                try:
                    from urllib.parse import urlparse, parse_qs
                    # Parse filters embedded in the endpoint URL
                    parsed = urlparse(endpoint)
                    filter_params = {
                        k: v[0]
                        for k, v in parse_qs(parsed.query).items()
                        if k not in ("limit", "offset", "search")
                    }

                    result = _api_call(endpoint, {})
                    items = result.get("data", [])

                    # Client-side filter (API may ignore query params server-side)
                    def _item_matches(item: dict) -> bool:
                        for k, v in filter_params.items():
                            iv = item.get(k)
                            if v.lower() == "true":
                                if not iv:
                                    return False
                            elif v.lower() == "false":
                                if iv:
                                    return False
                            else:
                                if str(iv).lower() != v.lower():
                                    return False
                        return True

                    items = [i for i in items if _item_matches(i)]

                    if items:
                        if not col_fields:
                            col_fields = list(items[0].keys())
                        else:
                            # Validate: if the specified fields don't exist, auto-detect
                            if not any(f in items[0] for f in col_fields):
                                col_fields = list(items[0].keys())
                                col_headers = []
                        headers = col_headers or col_fields
                        rows = [headers] + [
                            [str(item.get(f, "")) for f in col_fields]
                            for item in items
                        ]
                        sec["rows"] = rows
                    else:
                        sec["text"] = (sec.get("text", "") + " (No data found.)").strip()
                except Exception as fe:
                    sec["text"] = (sec.get("text", "") + f" [data fetch error: {fe}]").strip()

        gen_fn, suffix = _FORMAT_GENERATORS[fmt]
        data = gen_fn(title, sections)
        slug = re.sub(r"[^a-z0-9]+", "_", title.lower())[:28]
        filename = f"{slug}_{uuid.uuid4().hex[:8]}{suffix}"
        (_FILES_DIR / filename).write_bytes(data)
        public_base = os.environ.get("PUBLIC_BASE_URL", "http://localhost:8000").rstrip("/")
        artifact_url = f"{public_base}/files/{filename}"
        return {
            "artifact_url": artifact_url,
            "filename": filename,
            "byte_size": len(data),
            "sections_written": len(sections),
            "summary": f"Generated {fmt.upper()} '{title}' ({len(sections)} sections, {len(data):,} bytes).",
        }
    except Exception as exc:
        return {"error": f"artifact generation failed: {exc}"}


# ── Cache ─────────────────────────────────────────────────────────────────────

_CACHE: dict[str, dict] = {}


# ── System prompt ─────────────────────────────────────────────────────────────

_SYSTEM = """You are the Al Dente Company Brain, an AI agent for Al Dente S.r.l. (Italian pasta maker).

Use ONLY the provided tools to gather data. Never invent, estimate, or guess data.

Data sources:
- call_api /crm/*   → customers, opportunities, orders, invoices
- call_api /erp/*   → production orders (LOT-YYYY-#### IDs live here, not in /crm/orders), inventory (supplier_id on raw materials), suppliers, BOM, shipments
- call_api /calls/* → call list; /calls/{id}/transcript?search=<keyword> for transcript content
- search_kb         → product specs (shelf life, allergens, SKU), policies, price list, customer requirements
- search_transcripts → COUNT calls mentioning a keyword across the FULL log (aggregate/count use only; never use to look up a specific call)
- crm_channel_breakdown → value/count of opportunities grouped by channel (GDO/distributor/horeca)
- generate_artifact → create pdf/docx/pptx/xlsx binary file (ONLY for explicit binary format requests)

Artifact rules:
- "HTML deck" / "markdown" / "slides" without explicit file format → produce HTML/markdown INLINE in answer; artifact_url stays null.
- "PDF" / "docx" / "Word" / "pptx" / "PowerPoint" / "xlsx" / "Excel" / "downloadable" → call generate_artifact directly. Use 'data_source' in sections (e.g. data_source="/erp/inventory?type=finished_good&below_min=true") instead of 'rows' — Python fetches the data automatically. This keeps your tool call small and fast.

Rules:
1. If an entity doesn't exist in the data, say so explicitly — never invent it. LOT-YYYY-#### IDs are production lots in /erp/production-orders, not CRM orders.
2. Financial metrics (profit margin, cost, markup) are NOT stored anywhere in the available sources — if asked, state clearly they are not available after verifying the entity exists.
2. All numeric aggregates in tool results are pre-computed in Python — use them directly.
3. For the latest call with a named customer: if the customer_id is given use /calls?customer_id=X directly; if only the name is given, first /crm/customers?search=<partial_name> to get the id, then /calls?customer_id=<id>. Always pick the record with the most recent date — do NOT add an outcome filter.
4. Multi-hop BOM chain: /erp/bom?sku=PAS-X → get RAW-SKU → /erp/inventory?search=RAW-SKU (returns supplier_id) → /erp/suppliers (no filter, returns all; match by id field to get name).
5. 'Open' opportunities = stage qualification OR negotiation (won/lost are closed). To count/sum open opps: make TWO calls — /crm/opportunities?stage=qualification[&customer_id=X] and /crm/opportunities?stage=negotiation[&customer_id=X] — then ADD the two sum_value_eur and the two total counts from those responses.
6. For "by channel" groupings use crm_channel_breakdown — it joins opportunities to customers and sums in Python.
7. Quality complaint transcripts: use outcome_filter=complaint_open in search_transcripts. NEVER use search_transcripts to look up a specific call — for a specific call use /calls?customer_id=X then /calls/{id}/transcript?search=.
8. Answer concisely and factually. Do NOT show chain-of-thought or reasoning steps.
9. When a question includes a customer_id (e.g. CUST-0137), use it directly — do NOT pre-verify by calling /crm/customers. When a question only names a customer without an ID, call /crm/customers?search=<name_partial> to find the customer_id, then proceed. Only declare a customer non-existent if the name search returns zero results AND the relevant endpoints (opportunities, orders, calls) also return zero results.
10. For generation tasks (HTML deck, report), fetch the required data first (opportunities, orders, calls) using call_api, then produce the artifact inline.
11. For binary file requests: after generate_artifact returns successfully, write a 1-2 sentence summary answer (what the file contains) — the artifact_url will be attached automatically."""


# ── Tool definitions ──────────────────────────────────────────────────────────

_TOOLS: list[dict] = [
    {
        "type": "function",
        "function": {
            "name": "call_api",
            "description": (
                "Call any Al Dente API endpoint. All pages are fetched automatically. "
                "Numeric fields are pre-summed in the result as sum_<field>.\n"
                "Endpoints (filters are exact-match, case-sensitive):\n"
                "  /crm/customers  [search=<name_partial>, channel=(GDO|distributor|horeca), status=(active|inactive|prospect)]  — NOTE: no customer_id filter; cannot look up by ID here\n"
                "  /crm/customers/{id}\n"
                "  /crm/opportunities  [customer_id, stage=(qualification|negotiation|won|lost), owner]\n"
                "  /crm/orders  [customer_id, status=(open|in_production|shipped|delivered|cancelled), from, to]\n"
                "  /crm/invoices  [customer_id, status=(unpaid|paid|overdue), order_id]\n"
                "  /calls  [customer_id, type=(sales|support), outcome=(complaint_open|follow_up|order_placed|resolved), from, to]\n"
                "  /calls/{id}\n"
                "  /calls/{id}/transcript  [search=<keyword>, speaker]  — always use search= to filter\n"
                "  /erp/production-orders  [customer_id, status=(planned|in_progress|done|blocked), sku, from, to]\n"
                "  /erp/inventory  [type=(finished_good|raw_material), below_min=true, search=<sku>]  — raw_material rows include supplier_id\n"
                "  /erp/suppliers  [search=<name>, category=(semolina|wheat|packaging|labels|ink|logistics)]  — use no filter to list all; each record has id and name\n"
                "  /erp/bom  [sku=<PAS-XXX-###>]\n"
                "  /erp/shipments  [customer_id, order_id, status=(in_transit|delivered|delayed)]"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "endpoint": {
                        "type": "string",
                        "description": "API path, e.g. /crm/opportunities or /calls/CALL-58020/transcript",
                    },
                    "params": {
                        "type": "object",
                        "description": "Query parameters as string values",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": ["endpoint"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_kb",
            "description": (
                "Search Al Dente knowledge base documents (35 markdown files). "
                "Contains: product spec sheets (shelf life, allergens, PAS-* SKU details), "
                "quality & returns policies, 2026 wholesale price list, customer supply requirements. "
                "SKU queries (e.g. PAS-SPA-500) return the exact spec sheet automatically. "
                "Returns full document text — never truncated."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Keywords or SKU, e.g. 'PAS-SPA-500 shelf life allergens' or 'returns quality policy'",
                    },
                    "doc_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional: retrieve specific docs by ID, e.g. ['DOC-001', 'DOC-015']",
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_transcripts",
            "description": (
                "Scan every call transcript for a search keyword and count matching calls. "
                "Uses limit=1 probes — efficient even over 80+ calls. "
                "For quality-complaint counts, set outcome_filter=complaint_open. "
                "Returns: total_calls_searched, match_count, matching_call_ids, sample_segments."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "search_term": {
                        "type": "string",
                        "description": "Phrase to find in transcripts, e.g. 'broken pasta'",
                    },
                    "outcome_filter": {
                        "type": "string",
                        "description": "Pre-filter calls by outcome: complaint_open | follow_up | order_placed | resolved",
                    },
                    "customer_id": {
                        "type": "string",
                        "description": "Restrict search to one customer's calls",
                    },
                },
                "required": ["search_term"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "crm_channel_breakdown",
            "description": (
                "Get opportunity totals (count + sum of value_eur) grouped by customer channel "
                "(GDO / distributor / horeca). Joins opportunities to customers in Python — "
                "use this for questions like 'total value by channel' or 'opportunities per channel'. "
                "Optionally filter by opportunity stage before grouping."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "stage": {
                        "type": "string",
                        "description": "Filter opportunities by stage: qualification | negotiation | won | lost",
                    },
                    "customer_id": {
                        "type": "string",
                        "description": "Restrict to one customer (optional)",
                    },
                },
            },
        },
    },
]


# ── API layer (pagination + pre-computed aggregates) ─────────────────────────

def _headers() -> dict[str, str]:
    return {"Authorization": f"Bearer {os.environ['MOCK_API_TOKEN']}"}


def _fetch_all(
    cli: httpx.Client,
    url: str,
    params: dict[str, str],
    data_key: str = "data",
) -> tuple[list, int]:
    """Fetch all pages for a list endpoint. Returns (all_items, total)."""
    rp = dict(params)
    rp.setdefault("limit", "200")
    rp["offset"] = "0"
    hdrs = _headers()

    r = cli.get(url, params=rp, headers=hdrs)
    r.raise_for_status()
    first = r.json()
    items: list = list(first.get(data_key, []))
    total: int = first["pagination"]["total"]

    while len(items) < total:
        rp["offset"] = str(len(items))
        rp2 = cli.get(url, params=rp, headers=hdrs)
        rp2.raise_for_status()
        pg = rp2.json()
        new = pg.get(data_key, [])
        if not new:
            break
        items.extend(new)

    return items, total


def _api_call(endpoint: str, params: dict | None = None) -> dict:
    """Fully-paginated API call. Pre-computes sums for numeric fields."""
    base = os.environ["MOCK_API_BASE_URL"].rstrip("/")
    url = f"{base}{endpoint}"
    rp = {k: str(v) for k, v in (params or {}).items()}

    with httpx.Client(timeout=20.0) as cli:
        hdrs = _headers()
        rp2 = dict(rp)
        rp2.setdefault("limit", "200")
        rp2["offset"] = "0"
        r = cli.get(url, params=rp2, headers=hdrs)
        r.raise_for_status()
        first = r.json()

        # Transcript endpoint: { segments, pagination, call_id }
        if "segments" in first:
            items: list = list(first["segments"])
            total: int = first["pagination"]["total"]
            while len(items) < total:
                rp2["offset"] = str(len(items))
                rp3 = cli.get(url, params=rp2, headers=hdrs)
                rp3.raise_for_status()
                pg = rp3.json()
                new = pg.get("segments", [])
                if not new:
                    break
                items.extend(new)
            return {"segments": items, "total": total, "call_id": first.get("call_id")}

        # Standard list endpoint: { data, pagination }
        if "data" in first:
            items = list(first["data"])
            total = first["pagination"]["total"]
            while len(items) < total:
                rp2["offset"] = str(len(items))
                rp3 = cli.get(url, params=rp2, headers=hdrs)
                rp3.raise_for_status()
                pg = rp3.json()
                new = pg.get("data", [])
                if not new:
                    break
                items.extend(new)

            result: dict = {"data": items, "total": total}
            # Pre-compute numeric sums so LLM never has to do arithmetic
            if items and isinstance(items[0], dict):
                for key in items[0]:
                    if any(t in key.lower() for t in ("value", "amount", "price", "total", "quantity", "qty")):
                        try:
                            result[f"sum_{key}"] = sum(float(it.get(key) or 0) for it in items)
                        except (TypeError, ValueError):
                            pass
            return result

        # Single-object endpoint
        return first


# ── crm_channel_breakdown ─────────────────────────────────────────────────────

def _crm_channel_breakdown_impl(
    stage: str | None = None,
    customer_id: str | None = None,
) -> dict:
    """Python join: opportunities × customers → totals per channel. No LLM math."""
    base = os.environ["MOCK_API_BASE_URL"].rstrip("/")

    opp_params: dict[str, str] = {"limit": "200"}
    if stage:
        opp_params["stage"] = stage
    if customer_id:
        opp_params["customer_id"] = customer_id

    with httpx.Client(timeout=30.0) as cli:
        opps, _ = _fetch_all(cli, f"{base}/crm/opportunities", opp_params)
        customers, _ = _fetch_all(cli, f"{base}/crm/customers", {"limit": "200"})

    channel_map: dict[str, str] = {c["id"]: c.get("channel", "unknown") for c in customers}

    groups: dict[str, dict] = {}
    for opp in opps:
        ch = channel_map.get(opp.get("customer_id", ""), "unknown")
        if ch not in groups:
            groups[ch] = {"count": 0, "total_value_eur": 0.0, "opportunity_ids": []}
        groups[ch]["count"] += 1
        groups[ch]["total_value_eur"] += float(opp.get("value_eur") or 0)
        groups[ch]["opportunity_ids"].append(opp.get("id"))

    return {
        "breakdown_by_channel": groups,
        "total_opportunities": len(opps),
        "stage_filter": stage,
        "grand_total_value_eur": sum(g["total_value_eur"] for g in groups.values()),
    }


# ── search_transcripts ────────────────────────────────────────────────────────

def _search_transcripts_impl(
    search_term: str,
    outcome_filter: str | None = None,
    customer_id: str | None = None,
) -> dict:
    """Scan all (filtered) call transcripts for search_term. limit=1 probes for efficiency."""
    base = os.environ["MOCK_API_BASE_URL"].rstrip("/")

    call_params: dict[str, str] = {"limit": "200"}
    if outcome_filter:
        call_params["outcome"] = outcome_filter
    if customer_id:
        call_params["customer_id"] = customer_id

    with httpx.Client(timeout=90.0) as cli:
        all_calls, _ = _fetch_all(cli, f"{base}/calls", call_params)

        match_count = 0
        matching_ids: list[str] = []
        sample_segments: list[dict] = []

        for call in all_calls:
            cid = call["id"]
            tr = cli.get(
                f"{base}/calls/{cid}/transcript",
                params={"search": search_term, "limit": "1"},
                headers=_headers(),
            )
            if tr.status_code != 200:
                continue
            tr_data = tr.json()
            # Require at least 2 matching segments to exclude incidental mentions
            # (e.g. "alongside broken pasta" in a policy explanation, or "not broken pasta").
            # Genuine complaints always have 3+ hits; single-hit matches are passing references.
            if tr_data["pagination"]["total"] >= 2:
                match_count += 1
                matching_ids.append(cid)
                if len(sample_segments) < 6:
                    # Fetch up to 2 segments for context
                    tr2 = cli.get(
                        f"{base}/calls/{cid}/transcript",
                        params={"search": search_term, "limit": "2"},
                        headers=_headers(),
                    )
                    if tr2.status_code == 200:
                        for seg in tr2.json().get("segments", [])[:2]:
                            sample_segments.append({"call_id": cid, **seg})

    return {
        "search_term": search_term,
        "outcome_filter": outcome_filter,
        "total_calls_searched": len(all_calls),
        "match_count": match_count,
        "matching_call_ids": matching_ids,
        "sample_segments": sample_segments,
    }


# ── generate_artifact tool definition ────────────────────────────────────────

_TOOLS.append(
    {
        "type": "function",
        "function": {
            "name": "generate_artifact",
            "description": (
                "Generate a downloadable binary file (pdf / docx / pptx / xlsx). "
                "Call this ONLY when the user explicitly requests one of those formats "
                "or uses words like 'downloadable', 'Excel sheet', 'Word document', 'PowerPoint'. "
                "Do NOT call for 'HTML deck' or 'markdown' — produce those inline in the answer. "
                "IMPORTANT: use 'data_source' in sections instead of 'rows' to avoid sending large data. "
                "Python will fetch the data automatically. "
                "Returns artifact_url — include it in your final answer summary."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "format": {
                        "type": "string",
                        "enum": ["pdf", "docx", "pptx", "xlsx"],
                        "description": "File format: pdf, docx, pptx, or xlsx",
                    },
                    "title": {
                        "type": "string",
                        "description": "Document / presentation title",
                    },
                    "sections": {
                        "type": "array",
                        "description": (
                            "Content sections. For xlsx: one sheet per section. "
                            "For pptx: one slide per section. Each section may have "
                            "'heading', 'text' (paragraph), and/or 'rows' (table: "
                            "first row = headers, subsequent rows = data)."
                        ),
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "text": {"type": "string"},
                                "data_source": {
                                    "type": "string",
                                    "description": (
                                        "API endpoint to fetch rows automatically, e.g. "
                                        "'/erp/inventory?type=finished_good&below_min=true'. "
                                        "Use this instead of 'rows' to avoid sending large data in the call."
                                    ),
                                },
                                "column_fields": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": (
                                        "API field names to extract. "
                                        "Inventory fields: sku, description, on_hand, unit, min_stock, below_min, location. "
                                        "Leave empty to auto-include all fields."
                                    ),
                                },
                                "column_headers": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "Display header labels matching column_fields order",
                                },
                                "rows": {
                                    "type": "array",
                                    "description": "Explicit table rows (first row = headers). Use data_source instead when possible.",
                                    "items": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                },
                            },
                            "required": ["heading"],
                        },
                    },
                },
                "required": ["format", "title", "sections"],
            },
        },
    }
)


# ── Tool dispatcher ───────────────────────────────────────────────────────────

def _run_tool(name: str, args: dict) -> tuple[str, list[str], str]:
    """Execute one tool call. Returns (result_json, new_sources, verticale_cat)."""

    if name == "call_api":
        endpoint: str = args.get("endpoint", "")
        params: dict = args.get("params", {})
        try:
            result = _api_call(endpoint, params)
            segs = [s for s in endpoint.split("/") if s]
            cat = segs[0] if segs and segs[0] in ("crm", "erp", "calls") else "unknown"
            return json.dumps(result), [endpoint.lstrip("/")], cat
        except httpx.HTTPStatusError as e:
            body = {
                "error": "not_found" if e.response.status_code == 404 else "api_error",
                "http_status": e.response.status_code,
                "endpoint": endpoint,
            }
            return json.dumps(body), [], "unknown"
        except Exception as exc:
            return json.dumps({"error": str(exc), "endpoint": endpoint}), [], "unknown"

    if name == "search_kb":
        try:
            docs = _search_kb_impl(args.get("query", ""), args.get("doc_ids"))
            sources = [d["id"] for d in docs]
            return json.dumps(docs), sources, "kb"
        except Exception as exc:
            return json.dumps({"error": str(exc)}), [], "kb"

    if name == "search_transcripts":
        try:
            result = _search_transcripts_impl(
                args.get("search_term", ""),
                args.get("outcome_filter"),
                args.get("customer_id"),
            )
            return json.dumps(result), ["calls/transcripts"], "calls"
        except Exception as exc:
            return json.dumps({"error": str(exc)}), [], "calls"

    if name == "crm_channel_breakdown":
        try:
            result = _crm_channel_breakdown_impl(
                args.get("stage"),
                args.get("customer_id"),
            )
            return json.dumps(result), ["crm/opportunities", "crm/customers"], "crm"
        except Exception as exc:
            return json.dumps({"error": str(exc)}), [], "crm"

    if name == "generate_artifact":
        try:
            result = _generate_artifact_impl(
                args.get("format", ""),
                args.get("title", "Artifact"),
                args.get("sections", []),
            )
            return json.dumps(result), [], "kb"
        except Exception as exc:
            return json.dumps({"error": str(exc)}), [], "kb"

    return json.dumps({"error": f"unknown tool: {name}"}), [], "unknown"


def _strip_preamble(text: str) -> str:
    """Remove chain-of-thought preamble lines that some models emit before the answer."""
    _PREAMBLE_RE = re.compile(
        r"^\s*(now\s+i\s+have\s+the\s+answer[\.\:]?|"
        r"let\s+me\s+(provide|give|summarize|now\s+answer)[\w\s]*[\.\:]?|"
        r"based\s+on\s+(the\s+)?(data|information|results?)[\w\s]*,?\s*|"
        r"so\s+the\s+(correct\s+)?(answer|price|conclusion)\s+is[\:\.]?\s*)",
        re.IGNORECASE,
    )
    lines = text.split("\n")
    # Drop leading lines that are pure preamble
    while lines and _PREAMBLE_RE.match(lines[0].strip()):
        lines.pop(0)
    return "\n".join(lines).strip()


# ── Verticale routing ─────────────────────────────────────────────────────────

def _pick_verticale(cats: list[str]) -> str:
    counts: dict[str, int] = {"crm": 0, "erp": 0, "calls": 0, "kb": 0}
    for c in cats:
        if c in counts:
            counts[c] += 1
    best = max(counts, key=counts.get)  # type: ignore[arg-type]
    return best if counts[best] > 0 else "kb"


# ── Binary format detector ───────────────────────────────────────────────────
_BINARY_FMT_RE = re.compile(
    r"\b(pdf|docx|xlsx|pptx|excel\b|word\s+doc|powerpoint|downloadable)\b",
    re.IGNORECASE,
)
_INLINE_FMT_RE = re.compile(
    r"\b(html\s+(deck|page|report|slides?)|markdown|slide\s+deck)\b",
    re.IGNORECASE,
)


def _detect_binary_format(text: str) -> str | None:
    """Return 'pdf', 'docx', 'pptx', 'xlsx' if an explicit binary format is requested."""
    t = text.lower()
    if "pdf" in t:
        return "pdf"
    if "docx" in t or "word doc" in t:
        return "docx"
    if "pptx" in t or "powerpoint" in t:
        return "pptx"
    if "xlsx" in t or "excel" in t:
        return "xlsx"
    if "downloadable" in t:
        return "pdf"  # default binary to pdf
    return None


# ── Fast binary artifact path ─────────────────────────────────────────────────

_PLAN_SYSTEM = (
    "You are a data extraction planner. "
    "Given a question, output ONLY valid JSON — no markdown, no commentary.\n"
    "Schema: "
    '{"title":"string","sections":['
    '{"heading":"string",'
    '"data_source":"API endpoint with query params, e.g. /erp/inventory?type=finished_good&below_min=true",'
    '"column_fields":["field1","field2",...],'
    '"column_headers":["Label1","Label2",...]}'
    "]}\n"
    "Available endpoints and key fields:\n"
    "  /erp/inventory [type=finished_good|raw_material, below_min=true]  "
    "fields: sku, description, on_hand, unit, min_stock, below_min\n"
    "  /erp/production-orders [status=planned|in_progress|done|blocked, sku]  "
    "fields: lot_id, sku, status, planned_qty, produced_qty, start_date, end_date\n"
    "  /erp/suppliers  fields: id, name, category, country\n"
    "  /crm/customers  fields: id, name, channel, status\n"
    "  /crm/opportunities [stage=qualification|negotiation|won|lost]  "
    "fields: id, customer_id, title, stage, value_eur\n"
    "  /crm/orders [status=open|in_production|shipped|delivered]  "
    "fields: id, customer_id, sku, qty, status, date\n"
    "Output only the JSON object."
)


def _run_binary_artifact(
    question: str, binary_fmt: str, llm: "OpenAI", model: str, t0: float
) -> dict | None:
    """
    One-shot fast path: plan the artifact structure with one LLM call (no tools),
    then fetch data and generate the file in Python.
    Returns the /ask response dict, or None if it should fall back to the main loop.
    """
    remaining = _REQUEST_BUDGET - (time.monotonic() - t0)
    if remaining < 8:
        return None
    try:
        plan_resp = llm.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": _PLAN_SYSTEM},
                {"role": "user", "content": f"Question: {question}\nFormat: {binary_fmt}"},
            ],
            timeout=min(15.0, remaining - 5),
        )
        raw = (plan_resp.choices[0].message.content or "").strip()
        # Strip markdown code fences if present
        raw = re.sub(r"^```[a-z]*\n?", "", raw).rstrip("`").strip()
        plan = json.loads(raw)
    except Exception:
        return None  # fall back to main loop

    title = plan.get("title", "Report")
    sections = plan.get("sections", [])
    if not sections:
        return None

    remaining2 = _REQUEST_BUDGET - (time.monotonic() - t0)
    if remaining2 < 4:
        return None

    result = _generate_artifact_impl(binary_fmt, title, sections)
    if "error" in result:
        return None

    # Quick summary answer — no extra LLM call needed
    n_rows = sum(
        len(s.get("rows", [])) - 1  # subtract header
        for s in sections
        if s.get("rows")
    )
    summary = f"Generated {binary_fmt.upper()} '{title}'"
    if n_rows > 0:
        summary += f" with {n_rows} data rows"
    summary += "."

    return {
        "answer": summary,
        "sources": ["erp/inventory"],  # best-effort
        "verticale": _pick_verticale(["erp"]),
        "artifact_url": result.get("artifact_url"),
    }


# ── Agent loop ────────────────────────────────────────────────────────────────

_MAX_STEPS = 12
_REQUEST_BUDGET = 26.0  # total seconds allowed for the whole agent run


def run_agent(question: str) -> dict:
    """Run the agent loop and return the /ask response dict."""
    if question in _CACHE:
        return _CACHE[question]

    llm = OpenAI(
        api_key=os.environ["LLM_API_KEY"],
        base_url=os.environ["LLM_BASE_URL"],
    )
    model = os.environ["MODEL"]
    t0 = time.monotonic()

    # Normalise: wrap command-style fragments so any model engages with them
    q = question.strip()
    if q and not any(q.lower().startswith(w) for w in (
        "what", "who", "when", "where", "why", "how", "is ", "are ", "does ", "do ",
        "can ", "could ", "would ", "should ", "has ", "have ", "did ", "was ", "were ",
        "list", "show", "find", "get", "give", "generate", "create", "make", "build",
        "count", "calculate", "compute", "across",
    )):
        q = f"Please research and answer: {q}"

    binary_fmt = _detect_binary_format(q) if not _INLINE_FMT_RE.search(q) else None

    # Fast path for binary artifacts: one LLM planning call + Python generation
    if binary_fmt:
        fast = _run_binary_artifact(question, binary_fmt, llm, model, t0)
        if fast is not None:
            _CACHE[question] = fast
            return fast
        # Fall back: add a strong hint to the normal agent loop
        q = (
            f"{q}\n\n"
            f"IMPORTANT: Call generate_artifact immediately. Use 'data_source' to point to the API — "
            f"Python fetches the rows automatically. Example call:\n"
            f'{{"format":"{binary_fmt}","title":"...","sections":['
            f'{{"heading":"Section","data_source":"/erp/inventory?type=finished_good&below_min=true",'
            f'"column_fields":["sku","description","on_hand","min_stock"]}}]}}\n'
            f"Do NOT write descriptive text in sections. Do NOT include any URL in your answer."
        )

    messages: list[dict] = [
        {"role": "system", "content": _SYSTEM},
        {"role": "user", "content": q},
    ]

    sources: list[str] = []
    cats: list[str] = []
    artifact_url: str | None = None

    for _step in range(_MAX_STEPS):
        remaining = _REQUEST_BUDGET - (time.monotonic() - t0)
        if remaining < 3:
            break
        step_timeout = min(20.0, remaining - 1.5)
        # For binary artifacts: force generate_artifact on step 0 so the model fills
        # in data_source directly — avoids multiple pre-fetch LLM turns
        if binary_fmt and _step == 0:
            tc_choice: str | dict = {
                "type": "function",
                "function": {"name": "generate_artifact"},
            }
        else:
            tc_choice = "auto"
        try:
            resp = llm.chat.completions.create(
                model=model,
                messages=messages,
                tools=_TOOLS,
                tool_choice=tc_choice,
                timeout=step_timeout,
            )
        except Exception as exc:
            return _error_response(str(exc), sources, cats)

        msg = resp.choices[0].message

        # Handle reasoning models that put text in reasoning_content
        content: str = msg.content or ""
        if not content:
            content = getattr(msg, "reasoning_content", None) or ""
        # Strip Qwen3 <think>…</think> chain-of-thought blocks
        content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()
        # Strip Mistral / other models' reasoning preambles that leak into content
        content = _strip_preamble(content)

        if not msg.tool_calls:
            # Clean up any hallucinated artifact URLs the model put in text
            clean = re.sub(r"\[?artifact_url\]?:\s*\S+", "", content, flags=re.IGNORECASE).strip()
            clean = re.sub(r"https?://\S+\.(pdf|docx|pptx|xlsx)\S*", "", clean, flags=re.IGNORECASE).strip()
            out = {
                "answer": clean or "No answer could be generated.",
                "sources": list(dict.fromkeys(sources)),
                "verticale": _pick_verticale(cats),
                "artifact_url": artifact_url,
            }
            _CACHE[question] = out
            return out

        # Append assistant turn with tool_calls
        messages.append(
            {
                "role": "assistant",
                "content": content,
                "tool_calls": [
                    {
                        "id": tc.id,
                        "type": tc.type,
                        "function": {
                            "name": tc.function.name,
                            "arguments": tc.function.arguments,
                        },
                    }
                    for tc in msg.tool_calls
                ],
            }
        )

        for tc in msg.tool_calls:
            try:
                args = json.loads(tc.function.arguments)
            except json.JSONDecodeError:
                args = {}

            result_json, new_sources, cat = _run_tool(tc.function.name, args)
            sources.extend(new_sources)
            if cat:
                cats.append(cat)
            # Track artifact_url produced by generate_artifact
            try:
                rd = json.loads(result_json)
                if isinstance(rd, dict) and "artifact_url" in rd:
                    artifact_url = rd["artifact_url"]
            except Exception:
                pass

            messages.append(
                {"role": "tool", "tool_call_id": tc.id, "content": result_json}
            )

    # Step/time limit hit — request a summary answer
    remaining = _REQUEST_BUDGET - (time.monotonic() - t0)
    if remaining > 3:
        try:
            final = llm.chat.completions.create(
                model=model,
                messages=messages
                + [{"role": "user", "content": "Summarise your findings and give a final answer now."}],
                timeout=min(18.0, remaining - 1),
            )
            content = re.sub(
                r"<think>.*?</think>",
                "",
                final.choices[0].message.content or "",
                flags=re.DOTALL,
            ).strip()
            content = _strip_preamble(content)
            content = content or "Step limit reached; answer may be incomplete."
        except Exception:
            content = "Unable to complete the answer within the step limit."
    else:
        content = "The request could not be completed within the time budget."

    out = {
        "answer": content,
        "sources": list(dict.fromkeys(sources)),
        "verticale": _pick_verticale(cats),
        "artifact_url": artifact_url,
    }
    _CACHE[question] = out
    return out


def _error_response(msg: str, sources: list[str], cats: list[str]) -> dict:
    return {
        "answer": f"I cannot answer right now due to a technical issue: {msg}",
        "sources": list(dict.fromkeys(sources)),
        "verticale": _pick_verticale(cats),
        "artifact_url": None,
    }
